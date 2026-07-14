#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Populates the capstone .pptx template with real content, charts and
screenshots produced by the earlier pipeline scripts."""
from copy import deepcopy
from pptx import Presentation
from pptx.util import Emu
from pptx.oxml.ns import qn
from PIL import Image

SRC = "/home/claude/spacex_capstone/pptx_work/template.pptx"
OUT = "/home/claude/spacex_capstone/pptx_work/spacex_capstone_completed.pptx"
IMG = "/home/claude/spacex_capstone/images"
GH = "https://github.com/massimo7719/spacex-capstone"

prs = Presentation(SRC)
slides = prs.slides


def set_bullets(shape, items):
    """Replace a shape's paragraphs with `items`, preserving the formatting
    of the shape's first existing run per paragraph (cloned as needed)."""
    tf = shape.text_frame
    txBody = tf._txBody
    existing_ps = txBody.findall(qn("a:p"))
    if not existing_ps:
        return
    template_p = deepcopy(existing_ps[0])
    for p in existing_ps:
        txBody.remove(p)
    for text in items:
        new_p = deepcopy(template_p)
        runs = new_p.findall(qn("a:r"))
        if runs:
            first_run = runs[0]
            for r in runs[1:]:
                new_p.remove(r)
            t = first_run.find(qn("a:t"))
            if t is None:
                t = first_run.makeelement(qn("a:t"), {})
                first_run.append(t)
            t.text = text
        txBody.append(new_p)


def set_runs_simple(shape, texts):
    """For shapes where paragraph count == len(texts): set first run text
    of each paragraph directly (keeps formatting, no XML surgery)."""
    paras = shape.text_frame.paragraphs
    for p, text in zip(paras, texts):
        if p.runs:
            p.runs[0].text = text
            for r in p.runs[1:]:
                r.text = ""
        else:
            p.text = text


def clear_text(shape):
    tf = shape.text_frame
    for p in tf.paragraphs:
        for r in p.runs:
            r.text = ""


def delete_shape(shape):
    el = shape._element
    el.getparent().remove(el)


def nb(n, name):
    return f"Notebook {n:02d} ({name}) - full repo link in Appendix, slide 46"


def add_image_fit(slide, img_path, left, top, max_width, max_height, valign="middle"):
    with Image.open(img_path) as im:
        w, h = im.size
    aspect = w / h
    box_aspect = max_width / max_height
    if aspect > box_aspect:
        width = max_width
        height = int(max_width / aspect)
    else:
        height = max_height
        width = int(max_height * aspect)
    left_final = left + (max_width - width) // 2
    if valign == "middle":
        top_final = top + (max_height - height) // 2
    else:
        top_final = top
    return slide.shapes.add_picture(img_path, left_final, top_final, width=width, height=height)


def get_shape(slide, shape_id=None, name=None):
    for shape in slide.shapes:
        if shape_id is not None and shape.shape_id == shape_id:
            return shape
        if name is not None and shape.name == name:
            return shape
    return None


S = slides  # 0-indexed; slide N in the dump == S[N-1]

# ---------------------------------------------------------------- Slide 1
sl = S[0]
tb = get_shape(sl, shape_id=6)
tb.text_frame.paragraphs[0].runs[0].text = "Massimo"
tb.text_frame.paragraphs[1].runs[0].text = "July 14, 2026"

# ---------------------------------------------------------------- Slide 3 (Executive Summary)
sl = S[2]
set_runs_simple(get_shape(sl, shape_id=10), [
    "Methodology: 663 real Falcon 9 launches (June 2010 - July 2026) collected via web scraping of Wikipedia's "
    "launch-history pages, cleaned and analyzed with SQL and visualization, mapped and dashboarded, and used to "
    "train 4 classification models predicting first-stage landing success.",
    "Results: landing success rose from 0% (2010-2014) to 99%+ since 2021; the best model (Decision Tree, "
    "tuned via GridSearchCV) reached 96.2% test accuracy predicting landing outcome.",
])

# ---------------------------------------------------------------- Slide 4 (Introduction)
sl = S[3]
set_runs_simple(get_shape(sl, shape_id=5), [
    "SpaceX advertises Falcon 9 launches at $62 million, well below competitors' $165 million+, largely because "
    "it reuses the first-stage booster instead of expending it - reuse is only possible if the booster lands.",
    "Can historical mission parameters (flight number, payload mass, launch site, orbit) predict whether a "
    "Falcon 9 first-stage landing will succeed - the key cost driver for a competing bidder?",
])

# ---------------------------------------------------------------- Slide 6 (Methodology)
sl = S[5]
set_runs_simple(get_shape(sl, shape_id=7), [
    "Overview:",
    "Data collection:",
    "Web-scraped Wikipedia's Falcon 9/Falcon Heavy launch tables (663 flights); SpaceX REST API reference code "
    "included but unreachable from this sandbox",
    "Data wrangling:",
    "Parsed dates, payload mass, launch site, orbit and landing outcome/type; removed footnote artifacts; "
    "filtered to Falcon 9 only",
    "Exploratory data analysis: 6 visualizations + 10 SQL queries over the cleaned dataset",
    "Interactive analytics: static map and dashboard equivalents (Folium/Plotly Dash unavailable in the sandbox)",
    "Predictive analysis: 4 classification models trained and compared",
    "Models tuned via 5-fold GridSearchCV; evaluated with test-set accuracy and confusion matrix",
])

# ---------------------------------------------------------------- Slide 7 (Data Collection overview)
sl = S[6]
set_runs_simple(get_shape(sl, shape_id=5), [
    "Two data sources: SpaceX's public REST API (api.spacexdata.com) and Wikipedia's Falcon 9/Falcon Heavy "
    "launch-history tables.",
    "This sandbox's network could not reach the SpaceX API (blocked at the network layer), so the final dataset "
    "(663 real launches, 2010-2026) was built from the Wikipedia scrape; both approaches' code are in the "
    f"linked GitHub notebooks ({GH}).",
])

# ---------------------------------------------------------------- Slide 8 (Data Collection - SpaceX API)
sl = S[7]
set_bullets(get_shape(sl, shape_id=3), [
    "Reference implementation: GET api.spacexdata.com/v4/launches/past, joined with /v4/launchpads and "
    "/v4/payloads for site and mass detail.",
    "Not reachable from this session's sandboxed network (persistent Cloudflare 525 error despite the API's own "
    f"status page reporting 100% uptime). {nb(1, 'data_collection_api')}",
])
ph8 = get_shape(sl, shape_id=5)
ph8_pos = (ph8.left, ph8.top, ph8.width, ph8.height)
delete_shape(ph8)
add_image_fit(sl, f"{IMG}/15_flowchart_api.png", *ph8_pos)

# ---------------------------------------------------------------- Slide 9 (Data Collection - Scraping)
sl = S[8]
set_bullets(get_shape(sl, shape_id=3), [
    "Scraped Wikipedia's Falcon 9 launch tables (main page + 4 year-range archives) with requests + "
    "pandas.read_html; combined, cleaned, kept Falcon 9 only.",
    "663 real launches, flight 1 (2010-06-04) to 663 (2026-07-11).",
    nb(2, "scraping"),
])
ph9 = get_shape(sl, shape_id=2)
ph9_pos = (ph9.left, ph9.top, ph9.width, ph9.height)
delete_shape(ph9)
add_image_fit(sl, f"{IMG}/16_flowchart_scraping.png", *ph9_pos)

# ---------------------------------------------------------------- Slide 10 (Data Wrangling)
sl = S[9]
set_bullets(get_shape(sl, shape_id=5), [
    "Parsed launch dates, payload mass (kg), launch site (mapped to CCAFS SLC-40 / KSC LC-39A / VAFB SLC-4E "
    "with coordinates), orbit, customer, and landing outcome (success/failure + type).",
    "Removed Wikipedia footnote artifacts and duplicated note rows; filtered to Falcon 9 only; deduplicated by "
    "flight number. Final dataset: 663 rows, 47 missing payload-mass values (rideshare missions).",
    nb(3, "data_wrangling"),
])

# ---------------------------------------------------------------- Slide 11 (EDA visualization)
sl = S[10]
set_bullets(get_shape(sl, shape_id=5), [
    "6 charts: Flight Number vs. Launch Site, Payload vs. Launch Site, success rate by Orbit, Flight Number vs. "
    "Orbit, Payload vs. Orbit, and yearly success-rate trend - each colored by landing outcome.",
    nb(4, "eda_visualization"),
])

# ---------------------------------------------------------------- Slide 12 (EDA SQL)
sl = S[11]
set_bullets(get_shape(sl, shape_id=5), [
    "10 SQL queries on SPACEXTABLE (SQLite): unique launch sites; CCA-prefixed records; total NASA payload "
    "mass; average F9 v1.1 payload; first successful ground-pad landing date; drone-ship landings with "
    "4000-6000 kg payload; success/failure totals; max-payload boosters; failed 2015 drone-ship landings; "
    "ranked landing outcomes 2010-2017.",
    nb(5, "eda_sql"),
])

# ---------------------------------------------------------------- Slide 13 (Folium map)
sl = S[12]
set_bullets(get_shape(sl, shape_id=5), [
    "Site markers for all 3 launch complexes; per-launch outcome markers (green=success, red=failure) around "
    "each site; a proximity diagram for CCAFS SLC-40 (distance to coastline, highway, railway).",
    "Markers show where SpaceX operates from; outcome colors reveal both sites improved reliability over time; "
    "the proximity view shows why coastal sites were chosen - close to the ocean for downrange recovery, clear "
    "of population centers for safety.",
    f"{nb(6, 'folium_map')} (folium unavailable in this sandbox - static equivalent used, see slides 35-37)",
])

# ---------------------------------------------------------------- Slide 14 (Dashboard)
sl = S[13]
set_bullets(get_shape(sl, shape_id=5), [
    "Pie chart of successful landings by site; pie chart of outcomes at the best-performing site; payload-vs-"
    "outcome scatter split into 3 payload-range panels (in place of an interactive range slider).",
    "These views isolate which site contributes most successes, quantify the best single site's reliability, "
    "and show whether heavier payloads reduce landing success.",
    f"{nb(7, 'dashboard')} (plotly/dash unavailable in this sandbox - static equivalent used, see slides 39-41)",
])

# ---------------------------------------------------------------- Slide 15 (Predictive analysis)
sl = S[14]
set_bullets(get_shape(sl, shape_id=5), [
    "Trained Logistic Regression, SVM, Decision Tree, and KNN on flight number, payload mass, launch site, and "
    "orbit (one-hot encoded), each tuned via 5-fold GridSearchCV.",
    "Data split 80/20 (stratified), features standardized; best model selected by test-set accuracy: Decision "
    "Tree, 96.2%.",
    nb(8, "predictive_analysis"),
])

prs.save(OUT)
print(f"Salvato (parte 1): {OUT}")
