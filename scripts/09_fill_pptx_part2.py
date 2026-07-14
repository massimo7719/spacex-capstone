#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Part 2: slides 18-46 (EDA charts, SQL results, map/dashboard screenshots,
ML results, conclusions, appendix). Continues from the file part 1 saved."""
from copy import deepcopy
from pptx import Presentation
from pptx.util import Emu
from pptx.oxml.ns import qn
from PIL import Image

FILE = "/home/claude/spacex_capstone/pptx_work/spacex_capstone_completed.pptx"
IMG = "/home/claude/spacex_capstone/images"
GH = "https://github.com/massimo7719/spacex-capstone"

prs = Presentation(FILE)
S = prs.slides


def set_bullets(shape, items):
    tf = shape.text_frame
    txBody = tf._txBody
    existing_ps = txBody.findall(qn("a:p"))
    if not existing_ps:
        return
    template_p = deepcopy(existing_ps[0])
    for p in existing_ps:
        txBody.remove(p)
    for text in items:
        new_p = deepcopy(template_p)
        runs = new_p.findall(qn("a:r"))
        if runs:
            first_run = runs[0]
            for r in runs[1:]:
                new_p.remove(r)
            t = first_run.find(qn("a:t"))
            if t is None:
                t = first_run.makeelement(qn("a:t"), {})
                first_run.append(t)
            t.text = text
        txBody.append(new_p)


def set_runs_simple(shape, texts):
    paras = shape.text_frame.paragraphs
    for p, text in zip(paras, texts):
        if p.runs:
            p.runs[0].text = text
            for r in p.runs[1:]:
                r.text = ""
        else:
            p.text = text


def set_title(slide, text):
    for shape in slide.shapes:
        if shape.name.startswith("Title"):
            if shape.text_frame.paragraphs[0].runs:
                shape.text_frame.paragraphs[0].runs[0].text = text
                for r in shape.text_frame.paragraphs[0].runs[1:]:
                    r.text = ""
            else:
                shape.text_frame.text = text
            return


def clear_text(shape):
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            r.text = ""


def add_image_fit(slide, img_path, left, top, max_width, max_height, valign="middle"):
    with Image.open(img_path) as im:
        w, h = im.size
    aspect = w / h
    box_aspect = max_width / max_height
    if aspect > box_aspect:
        width = max_width
        height = int(max_width / aspect)
    else:
        height = max_height
        width = int(max_height * aspect)
    left_final = left + (max_width - width) // 2
    top_final = top + (max_height - height) // 2 if valign == "middle" else top
    return slide.shapes.add_picture(img_path, left_final, top_final, width=width, height=height)


def get_shape(slide, shape_id=None, name=None):
    for shape in slide.shapes:
        if shape_id is not None and shape.shape_id == shape_id:
            return shape
        if name is not None and shape.name == name:
            return shape
    return None


RIGHT_BOX = dict(left=Emu(4950000), top=Emu(1850000), width=Emu(6800000), height=Emu(4550000))

# ---------------------------------------------------------------- Slides 18-23 (EDA charts)
eda_slides = [
    (18, 3, "01_flightnumber_vs_site.png", [
        "Chart: flight number (x) vs. launch site (y), colored by landing outcome.",
        "Early flights at every site show more failures; from roughly flight 60 (~2016) onward, nearly every "
        "launch at all three sites lands successfully - landing technology matured independently of site.",
    ]),
    (19, 3, "02_payload_vs_site.png", [
        "Chart: payload mass in kg (x) vs. launch site (y), colored by landing outcome.",
        "VAFB SLC-4E and CCAFS SLC-40 handle payloads across the full range (mostly Starlink batches up to "
        "~17,500 kg); failures cluster in the early era rather than at any specific payload range.",
    ]),
    (20, 3, "03_success_rate_by_orbit.png", [
        "Chart: landing success rate (y) by orbit type (x), for orbits with 10+ launches.",
        "SSO (98.9%) and LEO (97.5%) have the highest success rates; GTO (76.2%) is lowest - GTO missions "
        "need more of the booster's fuel budget for the payload, leaving a harder landing profile.",
    ]),
    (21, 3, "04_flightnumber_vs_orbit.png", [
        "Chart: flight number (x) vs. orbit type (y), colored by landing outcome.",
        "GTO missions appear throughout the timeline including recent flights, and still show occasional "
        "failures even in the modern era - consistent with GTO being the hardest orbit to land from.",
    ]),
    (22, 3, "05_payload_vs_orbit.png", [
        "Chart: payload mass in kg (x) vs. orbit type (y), colored by landing outcome.",
        "LEO/SSO Starlink missions cluster at high payload mass (15,000-17,500 kg) with high success; GTO "
        "payloads are lighter on average but historically harder to land from.",
    ]),
    (23, 3, "06_yearly_success_trend.png", [
        "Chart: average yearly landing success rate, 2010-2026.",
        "0% in 2010-2014, rising through 62.5% (2016) and 77.8% (2017), a dip to 52.6% in 2018 (new Block 5 "
        "booster introduction), then a steady climb to 99%+ every year since 2021.",
    ]),
]
for slide_no, shape_id, img, texts in eda_slides:
    sl = S[slide_no - 1]
    set_bullets(get_shape(sl, shape_id=shape_id), texts)
    add_image_fit(sl, f"{IMG}/{img}", RIGHT_BOX["left"], RIGHT_BOX["top"], RIGHT_BOX["width"], RIGHT_BOX["height"])

# ---------------------------------------------------------------- Slides 24-33 (SQL results)
sql_slides = {
    24: ("All Launch Site Names", [
        "Query: SELECT DISTINCT LaunchSite FROM SPACEXTABLE ORDER BY LaunchSite;",
        "Result: CCAFS SLC-40, KSC LC-39A, VAFB SLC-4E - the three real-world Falcon 9 launch complexes.",
    ]),
    25: ("Launch Site Names Begin with 'CCA'", [
        "Query: SELECT ... WHERE LaunchSite LIKE 'CCA%' LIMIT 5;",
        "Result: first 5 CCAFS SLC-40 flights, 2010-2013 - Dragon Qualification Unit, COTS Demo 1 & 2, "
        "CRS-1, CRS-2, all early NASA cargo-resupply missions.",
    ]),
    26: ("Total Payload Mass", [
        "Query: SELECT Customer, SUM(PayloadMass) ... WHERE Customer LIKE '%NASA%' GROUP BY Customer;",
        "Result: NASA-related customer strings total from ~325 kg (rideshares) up to 129,050 kg for the "
        "'NASA (CTS)' commercial crew/cargo program - the single largest NASA payload category.",
    ]),
    27: ("Average Payload Mass by F9 v1.1", [
        "Query: SELECT AVG(PayloadMass) ... WHERE BoosterVersion LIKE '%v1.1%';",
        "Result: 2,534.7 kg average payload mass for the F9 v1.1 booster version (2013-2016 era, before "
        "Full Thrust and Block 5 upgrades).",
    ]),
    28: ("First Successful Ground Landing Date", [
        "Query: SELECT MIN(DateOnly) ... WHERE LandingType='Ground pad' AND LandingSuccess=1;",
        "Result: 2015-12-22 - the historic Flight 20 (OG2-2 mission) landing at Landing Zone 1, the first "
        "-ever successful orbital-class booster landing.",
    ]),
    29: ("Successful Drone Ship Landing with Payload between 4000 and 6000", [
        "Query: SELECT DISTINCT BoosterVersion ... WHERE LandingType='Drone ship' AND LandingSuccess=1 "
        "AND PayloadMass BETWEEN 4000 AND 6000;",
        "Result: 33 distinct booster/flight combinations, spanning early F9 FT boosters (B1021-2, B1022, "
        "B1026) through modern F9 B5 boosters (B1096-5).",
    ]),
    30: ("Total Number of Successful and Failure Mission Outcomes", [
        "Query: SELECT LaunchOutcome, COUNT(*) ... GROUP BY LaunchOutcome;",
        "Result: 661 mission successes vs. 2 mission failures - the overall Falcon 9 mission (not landing) "
        "success rate is 99.7%.",
    ]),
    31: ("Boosters Carried Maximum Payload", [
        "Query: SELECT BoosterVersion, PayloadMass ... WHERE PayloadMass = (SELECT MAX(PayloadMass) ...);",
        "Result: 20 different F9 B5 boosters tied at the maximum recorded payload mass of 17,500 kg - all "
        "recent high-capacity Starlink v2 mini deployments.",
    ]),
    32: ("2015 Launch Records", [
        "Query: SELECT ... WHERE LandingType='Drone ship' AND LandingSuccess=0 AND DateOnly LIKE '2015%';",
        "Result: 3 failed drone-ship landings in 2015 (flights 14, 17, 19; boosters B1012, B1015, B1018), "
        "all from CCAFS SLC-40, during the early F9 v1.1 landing-development era.",
    ]),
    33: ("Rank Landing Outcomes Between 2010-06-04 and 2017-03-20", [
        "Query: SELECT LandingOutcome, COUNT(*) ... GROUP BY LandingOutcome ORDER BY COUNT(*) DESC;",
        "Result (descending): Failure/Other 9, Failure/Drone ship 6, Success/Drone ship 5, Failure/Ocean 5, "
        "Success/Ground pad 3, Failure/Parachute 2, No attempt 1 - illustrating the trial-and-error phase "
        "before landings became routine.",
    ]),
}
for slide_no, (_, texts) in sql_slides.items():
    sl = S[slide_no - 1]
    set_bullets(get_shape(sl, shape_id=5), texts)

# ---------------------------------------------------------------- Slides 35-37 (Folium map screenshots)
map_slides = [
    (35, "Global Launch Site Locations", "07_map_launch_sites.png", [
        "All 3 Falcon 9 launch complexes plotted on a world map: CCAFS SLC-40 and KSC LC-39A on Florida's "
        "Space Coast, VAFB SLC-4E on California's central coast.",
        "Both US coasts are used so launches can fly out over open ocean in the required direction (east "
        "for equatorial/ISS orbits from Florida, south for polar/sun-synchronous orbits from California).",
        "Static equivalent of a Folium map (folium unavailable in this sandbox) - built with matplotlib "
        "and a public world-borders GeoJSON.",
    ]),
    (36, "Launch Outcomes by Site (Color-Coded)", "08_map_outcomes_colored.png", [
        "Per-launch outcome markers around each site (green=success, red=failure, gray=no attempt); points "
        "are randomly jittered around the true site coordinates to approximate a Folium MarkerCluster.",
        "Green dominates at both coasts, but red (failure) markers are visibly denser than average - a "
        "reminder that all 3 sites share the same overall improvement trend over time.",
        "Static equivalent of a Folium color-coded map (folium unavailable in this sandbox).",
    ]),
    (37, "Proximity Analysis - CCAFS SLC-40", "09_map_proximity.png", [
        "Distances from CCAFS SLC-40 to nearby geographic features, computed with the haversine formula "
        "from public landmark coordinates.",
        "The pad sits about 1 km from the Atlantic coastline (for downrange booster recovery and safety), "
        "and a few km from the nearest highway, railway, and town.",
        "Static equivalent of a Folium proximity map (folium unavailable in this sandbox).",
    ]),
]
for slide_no, title, img, texts in map_slides:
    sl = S[slide_no - 1]
    set_title(sl, title)
    ph = get_shape(sl, shape_id=5)
    orig_left, orig_top, orig_width, orig_height = ph.left, ph.top, ph.width, ph.height
    set_bullets(ph, texts)
    ph.height = Emu(1350000)
    img_top = orig_top + ph.height + Emu(60000)
    img_height = orig_height - ph.height - Emu(60000)
    add_image_fit(sl, f"{IMG}/{img}", orig_left, img_top, orig_width, img_height)

# ---------------------------------------------------------------- Slides 39-41 (Dashboard screenshots)
dash_slides = [
    (39, "Launch Success Share by Site", "10_dash_success_by_site.png", [
        "Pie chart of successful landings by launch site.",
        "CCAFS SLC-40 contributes the largest share of successful landings simply from launch volume (331 "
        "flights), even though its success rate (90.3%) is the lowest of the three sites.",
        "Static equivalent of a Plotly Dash pie chart (plotly/dash unavailable in this sandbox).",
    ]),
    (40, "Best-Performing Site: VAFB SLC-4E Outcomes", "11_dash_best_site_pie.png", [
        "Pie chart of outcomes at the site with the highest success ratio.",
        "VAFB SLC-4E has the best landing success rate of the three sites at 97.6% (212 attempts), just "
        "ahead of KSC LC-39A (96.4%) and well ahead of CCAFS SLC-40 (90.3%).",
        "Static equivalent of a Plotly Dash pie chart (plotly/dash unavailable in this sandbox).",
    ]),
    (41, "Payload vs. Launch Outcome by Payload Range", "12_dash_payload_vs_outcome.png", [
        "Payload-vs-outcome scatter split into three payload-range panels, standing in for an interactive "
        "range slider.",
        "Failures are not concentrated in any single payload band; the heaviest payloads (10,000-20,000 kg, "
        "mostly modern Starlink missions) actually show the highest success rate, reflecting how much "
        "landing reliability has improved as SpaceX has grown payload capacity.",
        "Static equivalent of an interactive Plotly Dash range-slider chart (plotly/dash unavailable here).",
    ]),
]
for slide_no, title, img, texts in dash_slides:
    sl = S[slide_no - 1]
    set_title(sl, title)
    ph = get_shape(sl, shape_id=5)
    orig_left, orig_top, orig_width, orig_height = ph.left, ph.top, ph.width, ph.height
    set_bullets(ph, texts)
    ph.height = Emu(1450000)
    img_top = orig_top + ph.height + Emu(60000)
    img_height = orig_height - ph.height - Emu(60000)
    add_image_fit(sl, f"{IMG}/{img}", orig_left, img_top, orig_width, img_height)

# ---------------------------------------------------------------- Slide 43 (Classification Accuracy)
sl = S[42]
ph43 = get_shape(sl, shape_id=5)
set_bullets(ph43, [
    "Bar chart comparing test-set accuracy of Logistic Regression, SVM, Decision Tree, and KNN, each tuned "
    "via 5-fold GridSearchCV.",
    "Decision Tree has the highest accuracy at 96.2% (max_depth=2); the other three models tie at 95.4%.",
])
img43_left = ph43.left + ph43.width + Emu(180000)
img43_width = Emu(12192000) - Emu(457200) - img43_left
add_image_fit(sl, f"{IMG}/13_ml_accuracy_comparison.png", img43_left, ph43.top, img43_width, ph43.height)

# ---------------------------------------------------------------- Slide 44 (Confusion Matrix)
sl = S[43]
ph44 = get_shape(sl, shape_id=5)
orig_left, orig_top, orig_width, orig_height = ph44.left, ph44.top, ph44.width, ph44.height
set_bullets(ph44, [
    "Confusion matrix of the best model (Decision Tree) on the held-out test set (131 flights): 123/123 "
    "successful landings correctly identified, but only 3/8 failures caught (5 failures misclassified as "
    "success) - the model is far more reliable at confirming a success than flagging a likely failure, a "
    "direct consequence of the 613-vs-41 class imbalance in the historical data.",
])
ph44.height = Emu(1500000)
img_top = orig_top + ph44.height + Emu(80000)
img_height = orig_height - ph44.height - Emu(80000)
add_image_fit(sl, f"{IMG}/14_ml_confusion_matrix.png", orig_left, img_top, orig_width, img_height)

# ---------------------------------------------------------------- Slide 45 (Conclusions)
sl = S[44]
set_bullets(get_shape(sl, shape_id=5), [
    "Falcon 9 first-stage landing success climbed from 0% (2010-2014) to 99%+ (every year since 2021), "
    "driven by the Block 5 booster and drone-ship/ground-pad recovery becoming routine.",
    "Launch site matters less than launch era: all three complexes converge on near-perfect reliability by "
    "the 2020s; VAFB SLC-4E currently has the single best track record (97.6%).",
    "GTO missions are the hardest to land from (76.2% success) versus LEO/SSO (97-99%), reflecting the "
    "smaller fuel margin left for the landing burn on high-energy orbits.",
    "A tuned Decision Tree predicts landing outcome with 96.2% test accuracy, but recall on the rare "
    "'failure' class is only 38% - useful for confirming likely successes, less so for flagging risk, given "
    "how imbalanced (613 vs. 41) the historical outcomes are.",
    "Both the SpaceX REST API and Folium/Plotly Dash labs could not be executed live in this sandboxed "
    "environment (network and package-installation restrictions); reference code for all three is included "
    f"in the project notebooks ({GH}) alongside the static equivalents actually used for this analysis.",
])

# ---------------------------------------------------------------- Slide 46 (Appendix)
sl = S[45]
set_bullets(get_shape(sl, shape_id=5), [
    f"Full project repository (data, scripts, notebooks, images): {GH}",
    "8 Jupyter notebooks: data collection (API reference + web scraping used), data wrangling, EDA "
    "visualization, EDA SQL, Folium map (reference + static equivalent), Plotly Dash dashboard (reference + "
    "static equivalent), predictive analysis (classification).",
    "Raw cleaned dataset: data/spacex_launches.csv (663 launches, 13 columns). SQL query results: "
    "data/sql_results.md. Model results: data/ml_results.json.",
    "All 16 charts/screenshots used in this presentation are in the images/ folder of the repository.",
])

prs.save(FILE)
print(f"Salvato (parte 2, completo): {FILE}")
